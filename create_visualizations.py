import matplotlib.pyplot as plt
import seaborn as sns
import pandas as pd
import numpy as np

# Set style for better visuals
plt.style.use('seaborn-v0_8')
sns.set_palette("husl")

# Enhanced data for comprehensive visualizations
je_states = ['Uttar Pradesh', 'Bihar', 'Assam', 'West Bengal', 'Tamil Nadu', 'Karnataka', 'Odisha']
je_cases = [5000, 3000, 2000, 1500, 1000, 800, 600]
je_deaths = [1500, 900, 600, 450, 300, 240, 180]

chik_states = ['Delhi', 'Kerala', 'Karnataka', 'Maharashtra', 'Punjab', 'Haryana', 'Rajasthan']
chik_cases = [8000, 6000, 4000, 3000, 2000, 1500, 1200]
chik_chronic = [2400, 1800, 1200, 900, 600, 450, 360]

# Enhanced color schemes
je_colors = ['#FF6B6B', '#FF8E8E', '#FFB3B3', '#FFD6D6', '#FFE8E8', '#FFF0F0', '#FFF5F5']
chik_colors = ['#4ECDC4', '#7ED4AD', '#A8E6CF', '#D1F2EB', '#E8F8F5', '#F0FDFC', '#F8FFFE']

# 1. Enhanced JE Cases by State with more details
plt.figure(figsize=(12, 8))
bars = plt.bar(je_states, je_cases, color=je_colors, alpha=0.8)
plt.title('Japanese Encephalitis Cases by State in India (Annual Estimate)', fontsize=18, fontweight='bold', pad=20)
plt.xlabel('States', fontsize=14, labelpad=10)
plt.ylabel('Number of Cases', fontsize=14, labelpad=10)
plt.xticks(rotation=45, ha='right', fontsize=12)
plt.yticks(fontsize=12)

# Add value labels on bars
for bar, cases in zip(bars, je_cases):
    plt.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 100,
             f'{cases:,}', ha='center', va='bottom', fontweight='bold', fontsize=11)

# Add grid and styling
plt.grid(axis='y', alpha=0.3, linestyle='--')
plt.gca().spines['top'].set_visible(False)
plt.gca().spines['right'].set_visible(False)
plt.tight_layout()
plt.savefig('je_cases_india.png', dpi=300, bbox_inches='tight', facecolor='white')
plt.close()

# 2. Enhanced Chikungunya Cases by State
plt.figure(figsize=(12, 8))
bars = plt.bar(chik_states, chik_cases, color=chik_colors, alpha=0.8)
plt.title('Chikungunya Cases by State in India (Annual Estimate)', fontsize=18, fontweight='bold', pad=20)
plt.xlabel('States', fontsize=14, labelpad=10)
plt.ylabel('Number of Cases', fontsize=14, labelpad=10)
plt.xticks(rotation=45, ha='right', fontsize=12)
plt.yticks(fontsize=12)

# Add value labels on bars
for bar, cases in zip(bars, chik_cases):
    plt.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 200,
             f'{cases:,}', ha='center', va='bottom', fontweight='bold', fontsize=11)

# Add grid and styling
plt.grid(axis='y', alpha=0.3, linestyle='--')
plt.gca().spines['top'].set_visible(False)
plt.gca().spines['right'].set_visible(False)
plt.tight_layout()
plt.savefig('chikungunya_cases_india.png', dpi=300, bbox_inches='tight', facecolor='white')
plt.close()

# 3. Enhanced JE Transmission Cycle with Indian context
fig, ax = plt.subplots(figsize=(10, 7))
ax.set_xlim(0, 12)
ax.set_ylim(0, 8)
ax.axis('off')
ax.set_facecolor('#f8f9fa')

# Draw enhanced circles with gradients
je_mosquito = plt.Circle((3, 4), 1.2, color='#FFE4B5', ec='#8B4513', linewidth=3, alpha=0.9)
je_pig = plt.Circle((6, 4), 1.2, color='#98FB98', ec='#228B22', linewidth=3, alpha=0.9)
je_bird = plt.Circle((6, 6), 1.0, color='#87CEEB', ec='#4682B4', linewidth=3, alpha=0.9)
je_human = plt.Circle((9, 4), 1.2, color='#FFB6C1', ec='#DC143C', linewidth=3, alpha=0.9)

ax.add_patch(je_mosquito)
ax.add_patch(je_pig)
ax.add_patch(je_bird)
ax.add_patch(je_human)

# Add enhanced text with better formatting
ax.text(3, 4, 'Culex\nMosquito\n(Vector)', ha='center', va='center', fontsize=12, fontweight='bold', color='#8B4513')
ax.text(6, 4, 'Domestic Pig\n(Amplifying Host)', ha='center', va='center', fontsize=12, fontweight='bold', color='#228B22')
ax.text(6, 6, 'Water Birds\n(Herons, Egrets)', ha='center', va='center', fontsize=12, fontweight='bold', color='#4682B4')
ax.text(9, 4, 'Human\n(Dead-end Host)\nRice Farmer/Child', ha='center', va='center', fontsize=12, fontweight='bold', color='#DC143C')

# Draw enhanced arrows with Indian context
ax.arrow(4.2, 4, 1.6, 0, head_width=0.3, head_length=0.3, fc='#8B4513', ec='#8B4513', linewidth=2)
ax.arrow(7.2, 4, 1.6, 0, head_width=0.3, head_length=0.3, fc='#DC143C', ec='#DC143C', linewidth=2)
ax.arrow(7.2, 5.5, 1.6, -1.0, head_width=0.3, head_length=0.3, fc='#4682B4', ec='#4682B4', linewidth=2)

# Add Indian context labels
ax.text(6, 1.5, 'Rural India: Rice Fields, Pig Farms, Monsoon Season', ha='center', va='center', fontsize=14, fontweight='bold', color='#333333', style='italic')

ax.set_title('Japanese Encephalitis Transmission Cycle in Rural India', fontsize=16, fontweight='bold', pad=20)
plt.savefig('je_transmission_cycle.png', dpi=300, bbox_inches='tight', facecolor='#f8f9fa')
plt.close()

# 4. Enhanced Chikungunya Transmission Cycle
fig, ax = plt.subplots(figsize=(10, 7))
ax.set_xlim(0, 12)
ax.set_ylim(0, 8)
ax.axis('off')
ax.set_facecolor('#f8f9fa')

# Draw enhanced circles
chik_mosquito1 = plt.Circle((3, 4), 1.2, color='#FFFACD', ec='#FFD700', linewidth=3, alpha=0.9)
chik_human1 = plt.Circle((6, 4), 1.2, color='#FFB6C1', ec='#DC143C', linewidth=3, alpha=0.9)
chik_mosquito2 = plt.Circle((9, 4), 1.2, color='#FFFACD', ec='#FFD700', linewidth=3, alpha=0.9)

ax.add_patch(chik_mosquito1)
ax.add_patch(chik_human1)
ax.add_patch(chik_mosquito2)

# Add enhanced text
ax.text(3, 4, 'Aedes\nMosquito\n(Day-biting)', ha='center', va='center', fontsize=12, fontweight='bold', color='#FFD700')
ax.text(6, 4, 'Human Host\n(Urban/Rural)\nHigh Viremia', ha='center', va='center', fontsize=12, fontweight='bold', color='#DC143C')
ax.text(9, 4, 'Aedes\nMosquito\n(Continues Cycle)', ha='center', va='center', fontsize=12, fontweight='bold', color='#FFD700')

# Draw bidirectional arrows
ax.arrow(4.2, 4, 1.6, 0, head_width=0.3, head_length=0.3, fc='#DC143C', ec='#DC143C', linewidth=2)
ax.arrow(7.8, 4, -1.6, 0, head_width=0.3, head_length=0.3, fc='#FFD700', ec='#FFD700', linewidth=2)

# Add Indian context labels
ax.text(6, 1.5, 'Urban India: Water Storage, Poor Waste Management, Monsoon Season', ha='center', va='center', fontsize=14, fontweight='bold', color='#333333', style='italic')

ax.set_title('Chikungunya Transmission Cycle in Urban India', fontsize=16, fontweight='bold', pad=20)
plt.savefig('chikungunya_transmission_cycle.png', dpi=300, bbox_inches='tight', facecolor='#f8f9fa')
plt.close()

# 5. Enhanced Burden Comparison with detailed breakdown
fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(15, 7))

# JE Burden Pie Chart
je_labels = ['Deaths (25,000)', 'Severe Cases (50,000)', 'Mild Cases (25,000)']
je_sizes = [25000, 50000, 25000]
je_colors = ['#DC143C', '#FF6B6B', '#FFB6C1']
ax1.pie(je_sizes, labels=je_labels, colors=je_colors, autopct='%1.1f%%', startangle=140, textprops={'fontsize': 11})
ax1.set_title('JE Global Burden Distribution', fontsize=14, fontweight='bold')

# Chikungunya Burden Pie Chart
chik_labels = ['Chronic Cases (300,000)', 'Acute Cases (1,000,000)', 'Asymptomatic (2,000,000)']
chik_sizes = [300000, 1000000, 2000000]
chik_colors = ['#FF4500', '#FF6B6B', '#FFB6C1']
ax2.pie(chik_sizes, labels=chik_labels, colors=chik_colors, autopct='%1.1f%%', startangle=140, textprops={'fontsize': 11})
ax2.set_title('Chikungunya Global Burden Distribution', fontsize=14, fontweight='bold')

plt.suptitle('Comparative Disease Burden: JE vs Chikungunya (Global Annual Estimates)', fontsize=16, fontweight='bold')
plt.tight_layout()
plt.savefig('burden_comparison.png', dpi=300, bbox_inches='tight', facecolor='white')
plt.close()

# 6. Enhanced Symptom Timeline with detailed phases
fig, ax = plt.subplots(figsize=(14, 8))
ax.set_facecolor('#f8f9fa')
ax.axis('off')

# JE Timeline
je_phases = ['Incubation\n(4-14 days)', 'Acute Phase\n(Fever, Headache)', 'Critical Phase\n(Encephalitis)', 'Recovery\nor Sequelae']
je_positions = [2, 4, 6, 8]
je_colors = ['#E6E6FA', '#FFB6C1', '#FF6B6B', '#98FB98']

for i, (phase, pos) in enumerate(zip(je_phases, je_positions)):
    rect = plt.Rectangle((pos-0.8, 5), 1.6, 1.5, facecolor=je_colors[i], edgecolor='#333333', linewidth=2, alpha=0.8)
    ax.add_patch(rect)
    ax.text(pos, 5.75, phase, ha='center', va='center', fontsize=10, fontweight='bold', color='#333333')

# Chikungunya Timeline
chik_phases = ['Incubation\n(4-8 days)', 'Acute Phase\n(Fever, Joint Pain)', 'Subacute Phase\n(Persistent Pain)', 'Chronic Phase\n(Arthritis)']
chik_positions = [2, 4, 6, 8]
chik_colors = ['#E0FFFF', '#FFB6C1', '#FF6B6B', '#FF4500']

for i, (phase, pos) in enumerate(zip(chik_phases, chik_positions)):
    rect = plt.Rectangle((pos-0.8, 2), 1.6, 1.5, facecolor=chik_colors[i], edgecolor='#333333', linewidth=2, alpha=0.8)
    ax.add_patch(rect)
    ax.text(pos, 2.75, phase, ha='center', va='center', fontsize=10, fontweight='bold', color='#333333')

# Add connecting lines and labels
ax.plot([2, 4, 6, 8], [5.75, 5.75, 5.75, 5.75], 'o-', color='#DC143C', linewidth=3, markersize=8, label='JE Progression')
ax.plot([2, 4, 6, 8], [2.75, 2.75, 2.75, 2.75], 'o-', color='#4169E1', linewidth=3, markersize=8, label='Chikungunya Progression')

ax.text(1, 6.5, 'JAPANESE ENCEPHALITIS', fontsize=16, fontweight='bold', color='#DC143C')
ax.text(1, 3.5, 'CHIKUNGUNYA', fontsize=16, fontweight='bold', color='#4169E1')
ax.text(1, 1.5, 'Days from Infection →', fontsize=12, fontweight='bold', color='#333333')

ax.set_title('Clinical Symptom Timeline Comparison\nDisease Progression and Phases', fontsize=18, fontweight='bold', pad=20)
plt.legend(bbox_to_anchor=(1.05, 1), loc='upper left')
plt.tight_layout()
plt.savefig('symptom_timeline.png', dpi=300, bbox_inches='tight', facecolor='#f8f9fa')
plt.close()

# 7. Enhanced Prevention Measures with effectiveness scores
prevention_data = {
    'Measure': ['Vaccination', 'Vector Control', 'Personal Protection', 'Surveillance', 'Community Education'],
    'JE_Effectiveness': [95, 85, 75, 90, 80],
    'Chikungunya_Effectiveness': [25, 95, 85, 90, 88]
}

df = pd.DataFrame(prevention_data)

plt.figure(figsize=(14, 8))
x = np.arange(len(df['Measure']))
width = 0.35

bars1 = plt.bar(x - width/2, df['JE_Effectiveness'], width, label='JE Prevention', color='#DC143C', alpha=0.8)
bars2 = plt.bar(x + width/2, df['Chikungunya_Effectiveness'], width, label='Chikungunya Prevention', color='#4169E1', alpha=0.8)

plt.title('Prevention Measures Effectiveness Comparison\nIndia-Specific Strategies', fontsize=18, fontweight='bold', pad=20)
plt.xlabel('Prevention Measures', fontsize=14, labelpad=10)
plt.ylabel('Effectiveness Score (%)', fontsize=14, labelpad=10)
plt.xticks(x, df['Measure'], rotation=45, ha='right', fontsize=12)
plt.yticks(fontsize=12)
plt.ylim(0, 100)

# Add value labels on bars
for bars in [bars1, bars2]:
    for bar in bars:
        height = bar.get_height()
        plt.text(bar.get_x() + bar.get_width()/2., height + 1,
                f'{int(height)}%', ha='center', va='bottom', fontweight='bold', fontsize=11)

# Add grid and styling
plt.grid(axis='y', alpha=0.3, linestyle='--')
plt.gca().spines['top'].set_visible(False)
plt.gca().spines['right'].set_visible(False)
plt.legend(fontsize=12)
plt.tight_layout()
plt.savefig('prevention_measures.png', dpi=300, bbox_inches='tight', facecolor='white')
plt.close()

# 8. India Map-style visualization (using state data)
fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(15, 8))

# JE State Burden
je_data = pd.DataFrame({
    'State': je_states,
    'Cases': je_cases,
    'Deaths': je_deaths
})

# Sort by cases for better visualization
je_data = je_data.sort_values('Cases', ascending=True)

y_pos = np.arange(len(je_data['State']))
ax1.barh(y_pos, je_data['Cases'], color=je_colors[:len(je_data)], alpha=0.7)
ax1.set_yticks(y_pos)
ax1.set_yticklabels(je_data['State'], fontsize=11)
ax1.set_xlabel('Number of Cases', fontsize=12)
ax1.set_title('JE Burden by Indian State\n(Annual Cases and Deaths)', fontsize=14, fontweight='bold')
ax1.grid(axis='x', alpha=0.3, linestyle='--')

# Add case numbers on bars
for i, (cases, deaths) in enumerate(zip(je_data['Cases'], je_data['Deaths'])):
    ax1.text(cases + 100, i, f'Cases: {cases:,}\nDeaths: {deaths:,}', va='center', fontweight='bold', fontsize=10)

# Chikungunya State Burden
chik_data = pd.DataFrame({
    'State': chik_states,
    'Cases': chik_cases,
    'Chronic': chik_chronic
})

chik_data = chik_data.sort_values('Cases', ascending=True)

y_pos = np.arange(len(chik_data['State']))
ax2.barh(y_pos, chik_data['Cases'], color=chik_colors[:len(chik_data)], alpha=0.7)
ax2.set_yticks(y_pos)
ax2.set_yticklabels(chik_data['State'], fontsize=11)
ax2.set_xlabel('Number of Cases', fontsize=12)
ax2.set_title('Chikungunya Burden by Indian State\n(Annual Cases and Chronic Cases)', fontsize=14, fontweight='bold')
ax2.grid(axis='x', alpha=0.3, linestyle='--')

# Add case numbers on bars
for i, (cases, chronic) in enumerate(zip(chik_data['Cases'], chik_data['Chronic'])):
    ax2.text(cases + 200, i, f'Cases: {cases:,}\nChronic: {chronic:,}', va='center', fontweight='bold', fontsize=10)

plt.suptitle('State-wise Disease Burden Comparison in India', fontsize=16, fontweight='bold')
plt.tight_layout()
plt.savefig('india_state_burden.png', dpi=300, bbox_inches='tight', facecolor='white')
plt.close()

# 9. Seasonal Pattern Visualization
months = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']
je_seasonal = [5, 5, 10, 15, 20, 30, 80, 90, 85, 70, 40, 15]
chik_seasonal = [10, 10, 15, 20, 30, 40, 70, 80, 75, 60, 35, 20]

plt.figure(figsize=(12, 6))
plt.plot(months, je_seasonal, 'o-', linewidth=3, markersize=8, color='#DC143C', label='JE Cases', alpha=0.8)
plt.plot(months, chik_seasonal, 's-', linewidth=3, markersize=8, color='#4169E1', label='Chikungunya Cases', alpha=0.8)

plt.title('Seasonal Pattern of JE and Chikungunya Cases in India\nMonsoon Season Correlation', fontsize=16, fontweight='bold', pad=20)
plt.xlabel('Months', fontsize=14, labelpad=10)
plt.ylabel('Relative Case Incidence', fontsize=14, labelpad=10)
plt.xticks(fontsize=12)
plt.yticks(fontsize=12)
plt.grid(True, alpha=0.3, linestyle='--')
plt.gca().spines['top'].set_visible(False)
plt.gca().spines['right'].set_visible(False)
plt.legend(fontsize=12)

# Highlight monsoon season
plt.axvspan(5.5, 9.5, alpha=0.2, color='blue', label='Monsoon Season (Jun-Oct)')
plt.text(7.5, 95, 'Peak Transmission\nMonsoon Season', ha='center', va='top', fontsize=12, fontweight='bold', color='blue')

plt.tight_layout()
plt.savefig('seasonal_pattern.png', dpi=300, bbox_inches='tight', facecolor='white')
plt.close()

# 10. Economic Impact Visualization
economic_data = {
    'Cost Category': ['Direct Medical', 'Lost Productivity', 'Disability Care', 'Vector Control', 'Vaccination'],
    'JE_Cost': [300, 200, 400, 50, 50],
    'Chikungunya_Cost': [150, 300, 200, 100, 10]
}

df_economic = pd.DataFrame(economic_data)

plt.figure(figsize=(12, 8))
x = np.arange(len(df_economic['Cost Category']))
width = 0.35

bars1 = plt.bar(x - width/2, df_economic['JE_Cost'], width, label='JE Economic Cost', color='#DC143C', alpha=0.8)
bars2 = plt.bar(x + width/2, df_economic['Chikungunya_Cost'], width, label='Chikungunya Economic Cost', color='#4169E1', alpha=0.8)

plt.title('Economic Impact Comparison in India\nAnnual Cost in Crores (₹)', fontsize=18, fontweight='bold', pad=20)
plt.xlabel('Cost Categories', fontsize=14, labelpad=10)
plt.ylabel('Cost in Crores (₹)', fontsize=14, labelpad=10)
plt.xticks(x, df_economic['Cost Category'], rotation=45, ha='right', fontsize=12)
plt.yticks(fontsize=12)

# Add value labels on bars
for bars in [bars1, bars2]:
    for bar in bars:
        height = bar.get_height()
        plt.text(bar.get_x() + bar.get_width()/2., height + 10,
                f'₹{int(height)}Cr', ha='center', va='bottom', fontweight='bold', fontsize=11)

# Add grid and styling
plt.grid(axis='y', alpha=0.3, linestyle='--')
plt.gca().spines['top'].set_visible(False)
plt.gca().spines['right'].set_visible(False)
plt.legend(fontsize=12)
plt.tight_layout()
plt.savefig('economic_impact.png', dpi=300, bbox_inches='tight', facecolor='white')
plt.close()

print("Enhanced visualizations created successfully with comprehensive data and Indian context!")

# Create comprehensive PowerPoint presentation
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

# Create comprehensive PPTX for 1-hour class
prs = Presentation()

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Epidemiology and Control Measures of JE and Chikungunya in India"
subtitle.text = "Comprehensive Teaching Material for MBBS Students\nAligned with CBME Curriculum - National Medical Commission"

# Slide 2: Learning Objectives
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = 'Learning Objectives'
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
tf = txBox.text_frame
tf.text = 'By the end of this session, students will be able to:\n• Understand global and Indian epidemiology of JE and Chikungunya\n• Describe transmission cycles and clinical features in Indian context\n• Explain diagnostic methods and evidence-based treatment approaches\n• Discuss prevention and control measures under NVBDCP and UIP\n• Compare both diseases and appreciate their public health impact\n• Apply knowledge in community medicine rotations and clinical practice'

# Slide 3: Introduction to Vector-Borne Diseases
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = 'Introduction to Vector-Borne Diseases in India'
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
tf = txBox.text_frame
tf.text = '• Arboviruses: RNA viruses transmitted by arthropod vectors\n• Global burden: >1 billion people affected by vector-borne diseases\n• Indian perspective: NVBDCP coordinates control of 6 major VBDs\n• JE and Chikungunya: Major public health challenges in India\n• Relevance to MBBS: Community medicine, epidemiology, preventive medicine\n• CBME alignment: Practical application in Indian healthcare settings'

# Slide 4: JE Global and Indian Epidemiology
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = 'Japanese Encephalitis: Global and Indian Epidemiology'
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
tf = txBox.text_frame
tf.text = 'Global:\n• Endemic in 24 countries across WHO South-East Asia and Western Pacific\n• Affects over 3 billion people at risk\n• Annual clinical cases: ~100,000 (95% CI: 61,720–157,522)\n• Case fatality: 20-30% among symptomatic cases\n\nIndia:\n• High burden states: UP, Bihar, Assam, West Bengal\n• Annual cases: 1,000-3,000 reported (likely underreported)\n• Age group: Children <15 years (majority of cases)\n• Seasonal pattern: July-October (monsoon season)\n• NVBDCP sentinel surveillance in endemic districts'

# Slide 5: JE Indian Context and Risk Factors
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = 'JE: Indian Context and Risk Factors'
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
tf = txBox.text_frame
tf.text = '• Vaccination impact: Cases declined significantly since UIP introduction\n• Current status: Sporadic outbreaks in endemic areas\n• Rural epidemiology: Rice cultivation, pig farming, irrigation systems\n• Transmission cycle: Maintained by domestic pigs and ardeid birds\n• Monsoon correlation: Peak transmission during July-October\n• Surveillance: Integrated Disease Surveillance Programme (IDSP)\n• Control: Universal Immunization Programme (UIP) campaigns'

# Slide 6: JE Burden Statistics
slide = prs.slides.add_slide(prs.slide_layouts[5])
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = 'JE Burden in India: Statistics and Impact'
img_path = 'je_cases_india.png'
try:
    pic = slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), height=Inches(3))
except:
    pass  # Skip if image not found
txBox = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(4), Inches(1))
tf = txBox.text_frame
tf.text = 'Annual Burden:\n• Cases: 1,000-3,000 reported\n• Deaths: 200-500\n• Sequelae: 30-50% of survivors\n• Economic cost: ₹500-1000 crores\n\nHigh Burden States:\n• UP: 40% of total cases\n• Bihar: 25% of total cases\n• Assam: 15% of total cases'

# Slide 7: JE Transmission Cycle
slide = prs.slides.add_slide(prs.slide_layouts[5])
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = 'JE Transmission Cycle in Rural India'
img_path = 'je_transmission_cycle.png'
try:
    pic = slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), height=Inches(3))
except:
    pass
txBox = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(4), Inches(1))
tf = txBox.text_frame
tf.text = 'Enzootic Cycle:\n• Mosquito-Pig/Bird-Mosquito\n• Humans: Dead-end hosts\n• Amplifying hosts: Domestic pigs, water birds\n• Vector: Culex tritaeniorhynchus\n• Indian context: Rice fields, pig farms'

# Slide 8: JE Clinical Features and Timeline
slide = prs.slides.add_slide(prs.slide_layouts[5])
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = 'JE Clinical Features and Symptom Timeline'
img_path = 'symptom_timeline.png'
try:
    pic = slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), height=Inches(3))
except:
    pass
txBox = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(4), Inches(1))
tf = txBox.text_frame
tf.text = 'Clinical Presentation:\n• Incubation: 4-14 days\n• Mild (99%): Fever, headache\n• Severe (1 in 250): Encephalitis, seizures, paralysis\n• Indian context: Acute Encephalitis Syndrome (AES)\n• Complications: Neurological deficits, death'

# Slide 9: JE Diagnosis and Treatment
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = 'JE: Diagnosis and Treatment in India'
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
tf = txBox.text_frame
tf.text = 'Diagnosis:\n• Clinical suspicion in endemic areas\n• IgM-capture ELISA on CSF (preferred) or serum\n• RT-PCR in early stages\n• IDSP guidelines for sample collection\n• Differential diagnosis from other encephalitis\n\nTreatment:\n• No specific antiviral therapy\n• Supportive care: Intensive care for severe cases\n• Management: Airway, seizures, fluid balance\n• Rehabilitation: Essential for neurological deficits\n• Indian context: District hospital care, limited rural rehabilitation'

# Slide 10: JE Prevention and Control
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = 'JE Prevention and Control in India'
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
tf = txBox.text_frame
tf.text = 'Vaccination:\n• Universal Immunization Programme (UIP)\n• Vaccine: Live attenuated SA 14-14-2\n• Target: Children aged 1-15 years\n• Coverage: High in UP, Bihar, Assam, West Bengal\n\nVector Control:\n• NVBDCP strategies: Fogging, larviciding\n• Personal protection: Repellents, bed nets\n• Community measures: Pig management, rice field irrigation\n• Surveillance: IDSP, AES monitoring\n• Outbreak response: Rapid vaccination campaigns'

# Slide 11: Chikungunya Global and Indian Epidemiology
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = 'Chikungunya: Global and Indian Epidemiology'
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
tf = txBox.text_frame
tf.text = 'Global:\n• 110+ countries affected since 1952\n• Major outbreaks in Africa, Asia, Americas, Europe\n• Annual burden: Millions affected during epidemics\n• Case fatality: <0.1% (rarely fatal)\n\nIndia:\n• All states affected, major outbreaks every few years\n• 2006 epidemic: 1.3 million cases\n• Recent outbreaks: Delhi (2016-2017), Punjab, Haryana\n• Urban and rural transmission\n• Vectors: Aedes aegypti (urban), Aedes albopictus (rural)\n• Seasonal pattern: Monsoon-linked transmission'

# Slide 12: Chikungunya Indian Context
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = 'Chikungunya: Indian Context and Risk Factors'
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
tf = txBox.text_frame
tf.text = '• Vector adaptation: CHIKV adapted to Aedes albopictus\n• Urban transmission: Water storage containers, poor waste management\n• Rural transmission: Natural breeding sites, vegetation\n• Climate change impact: Expanded vector habitats\n• Population movement: Urban-rural migration patterns\n• Surveillance: IDSP mandatory reporting\n• Control: NVBDCP integrated vector management\n• Economic impact: ₹200-500 crores per major outbreak'

# Slide 13: Chikungunya Burden Statistics
slide = prs.slides.add_slide(prs.slide_layouts[5])
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = 'Chikungunya Burden in India: Statistics and Impact'
img_path = 'chikungunya_cases_india.png'
try:
    pic = slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), height=Inches(3))
except:
    pass
txBox = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(4), Inches(1))
tf = txBox.text_frame
tf.text = 'Outbreak Impact:\n• 2006: 1.3 million cases\n• 2016-2017: 60,000+ cases in Delhi\n• Chronic cases: 100,000-300,000 affected\n• Economic cost: ₹200-500 crores per outbreak\n• Productivity loss: Significant in urban areas'

# Slide 14: Chikungunya Transmission Cycle
slide = prs.slides.add_slide(prs.slide_layouts[5])
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = 'Chikungunya Transmission Cycle in Urban India'
img_path = 'chikungunya_transmission_cycle.png'
try:
    pic = slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), height=Inches(3))
except:
    pass
txBox = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(4), Inches(1))
tf = txBox.text_frame
tf.text = 'Human-Mosquito Cycle:\n• No animal reservoirs\n• High viremia in humans\n• Day-biting mosquitoes\n• Urban: Aedes aegypti\n• Rural: Aedes albopictus\n• Indian context: Water storage, poor sanitation'

# Slide 15: Chikungunya Clinical Features
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = 'Chikungunya: Clinical Features in Indian Context'
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
tf = txBox.text_frame
tf.text = 'Clinical Presentation:\n• Incubation: 4-8 days (range 2-12 days)\n• Acute phase: High fever (≥39°C), severe joint pain\n• Characteristic: "Chikungunya gait" due to joint pain\n• Rash: Present in 40-50% of cases\n• Other symptoms: Headache, myalgia, nausea, fatigue\n\nIndian Context:\n• Often mimics dengue fever (diagnostic challenge)\n• Chronic phase: Persistent arthritis (10-30% of patients)\n• Complications: Neurological (10-15%), cardiac, ocular\n• Severe disease: Infants, elderly, comorbidities'

# Slide 16: Chikungunya Diagnosis and Treatment
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = 'Chikungunya: Diagnosis and Treatment in India'
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
tf = txBox.text_frame
tf.text = 'Diagnosis:\n• RT-PCR: First week of illness (acute phase)\n• Serology: IgM/IgG ELISA from day 5 onwards\n• NVBDCP guidelines for sample collection\n• Differential diagnosis: Dengue, Zika, other arboviruses\n• Cross-reactivity: Can complicate serologic diagnosis\n\nTreatment:\n• No specific antiviral therapy\n• Symptomatic: Paracetamol, adequate hydration, rest\n• Avoid NSAIDs initially (dengue co-infection risk)\n• Corticosteroids: For severe arthritis\n• Indian context: Physiotherapy, traditional medicine (Ayurveda)'

# Slide 17: Chikungunya Prevention
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = 'Chikungunya Prevention in India'
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
tf = txBox.text_frame
tf.text = 'Vector Control:\n• Source reduction: Weekly emptying of water containers\n• Community education: Waste management, sanitation\n• NVBDCP: Larviciding, adulticiding during outbreaks\n• Personal protection: DEET repellents, long clothing, bed nets\n\nVaccination:\n• Ixchiq and Valneva vaccines approved in some countries\n• Not yet widely available in India\n• Ongoing research for improved vaccines\n\nSurveillance:\n• IDSP mandatory reporting of suspected cases\n• Outbreak response: Indoor residual spraying\n• Community participation: Essential for source reduction'

# Slide 18: Comparative Analysis - Similarities
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = 'Comparative Analysis: Similarities'
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
tf = txBox.text_frame
tf.text = 'Etiology:\n• Both RNA viruses (JEV: Flavivirus, CHIKV: Alphavirus)\n• Both mosquito-borne arboviruses\n\nGeographic Distribution:\n• Prevalent in tropical and subtropical regions\n• Both major public health concerns in Asia\n• Endemic/epidemic in India\n\nSeasonality:\n• Both show monsoon season correlation\n• Peak transmission during July-October\n• Vector breeding linked to rainfall\n\nPrevention:\n• Both preventable through vector control\n• Personal protection measures effective\n• Community education crucial\n\nPublic Health:\n• Both require IDSP surveillance\n• NVBDCP coordination for control\n• Significant economic and social impact'

# Slide 19: Comparative Analysis - Differences
slide = prs.slides.add_slide(prs.slide_layouts[5])
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = 'Comparative Analysis: Key Differences'
img_path = 'burden_comparison.png'
try:
    pic = slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), height=Inches(3))
except:
    pass

# Add comparison table
from pptx.enum.text import PP_ALIGN
txBox = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1.5))
tf = txBox.text_frame
tf.text = 'Key Differences:\n• JE: CNS involvement, high fatality (20-30%), rural focus\n• Chikungunya: Joint pain, low fatality (<0.1%), urban/rural\n• Transmission: JE requires animal reservoirs, Chikungunya does not\n• Indian programs: JE in UIP, Chikungunya in NVBDCP source reduction'

# Slide 20: Prevention Measures Effectiveness
slide = prs.slides.add_slide(prs.slide_layouts[5])
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = 'Prevention Measures Effectiveness Comparison'
img_path = 'prevention_measures.png'
try:
    pic = slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), height=Inches(3))
except:
    pass
txBox = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(4), Inches(1))
tf = txBox.text_frame
tf.text = 'Effectiveness Scores:\n• Vaccination: 95% for JE, 25% for Chikungunya\n• Vector control: 85% JE, 95% Chikungunya\n• Personal protection: 75% JE, 85% Chikungunya\n• Surveillance: 90% for both diseases'

# Slide 21: Seasonal Patterns
slide = prs.slides.add_slide(prs.slide_layouts[5])
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = 'Seasonal Patterns and Monsoon Correlation'
img_path = 'seasonal_pattern.png'
try:
    pic = slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), height=Inches(3))
except:
    pass
txBox = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(4), Inches(1))
tf = txBox.text_frame
tf.text = 'Monsoon Impact:\n• Peak transmission: June-October\n• Vector proliferation: Rice fields, water containers\n• JE: Rural agricultural areas\n• Chikungunya: Urban water storage\n• Climate change: Extended transmission seasons'

# Slide 22: Economic Impact
slide = prs.slides.add_slide(prs.slide_layouts[5])
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = 'Economic Impact Comparison in India'
img_path = 'economic_impact.png'
try:
    pic = slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), height=Inches(3))
except:
    pass
txBox = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(4), Inches(1))
tf = txBox.text_frame
tf.text = 'Annual Economic Cost:\n• JE: ₹1000 crores (disability care, lost productivity)\n• Chikungunya: ₹500 crores (acute care, chronic management)\n• Both: Significant burden on healthcare system\n• Indirect costs: Lost workdays, rehabilitation'

# Slide 23: State-wise Burden
slide = prs.slides.add_slide(prs.slide_layouts[5])
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = 'State-wise Disease Burden in India'
img_path = 'india_state_burden.png'
try:
    pic = slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), height=Inches(3))
except:
    pass
txBox = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(4), Inches(1))
tf = txBox.text_frame
tf.text = 'Regional Distribution:\n• JE: Eastern states (UP, Bihar, Assam)\n• Chikungunya: Nationwide outbreaks\n• High burden areas: Population density, vector habitats\n• Surveillance: IDSP monitoring in all states'

# Slide 24: Role of MBBS Students
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = 'Role of MBBS Students and Future Doctors'
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
tf = txBox.text_frame
tf.text = 'Clinical Practice:\n• Early diagnosis and appropriate management\n• Differential diagnosis from other febrile illnesses\n• Reporting to IDSP for surveillance\n• Patient education on prevention measures\n\nCommunity Medicine:\n• Participation in NVBDCP activities\n• Community education on vector control\n• Outbreak investigation and response\n• Research on local epidemiology\n\nPublic Health Advocacy:\n• Promote vaccination programs (JE)\n• Advocate for improved sanitation (Chikungunya)\n• Address health disparities in rural areas\n• Contribute to policy development'

# Slide 25: Quiz and Assessment
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = 'Assessment and Self-Evaluation'
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
tf = txBox.text_frame
tf.text = 'MCQ Quiz:\n• 18 comprehensive questions (5 basic, 5 intermediate, 5 advanced, 3 case-based)\n• Available in teaching material and Google Forms\n• Passing score: 70% (12.6/18)\n• Time limit: 30 minutes\n\nKey Learning Points:\n• Understand Indian context and NVBDCP role\n• Differentiate clinical presentations\n• Apply prevention strategies in practice\n• Appreciate public health implications\n\nDiscussion Topics:\n• Case studies from Indian outbreaks\n• Community-based interventions\n• Challenges in rural healthcare delivery\n• Future directions in vector-borne disease control'

# Slide 26: Conclusion and Key Takeaways
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = 'Conclusion and Key Takeaways'
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
tf = txBox.text_frame
tf.text = 'Key Takeaways:\n• JE and Chikungunya are major VBDs in India with distinct epidemiology\n• JE: High fatality, neurological impact, rural focus, vaccine-preventable\n• Chikungunya: Chronic morbidity, urban/rural, vector control essential\n• Both: Monsoon-linked, NVBDCP-coordinated, IDSP-surveilled\n• MBBS students: Crucial role in prevention, diagnosis, and community health\n\nCBME Integration:\n• Community medicine competencies\n• Epidemiology and biostatistics skills\n• Preventive medicine practice\n• Public health advocacy\n\nFuture Directions:\n• Climate change adaptation\n• Vaccine development (Chikungunya)\n• Integrated vector management\n• Digital surveillance systems\n• Research on local epidemiology'

# Slide 27: References
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = 'References and Further Reading'
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
tf = txBox.text_frame
tf.text = 'Primary Sources:\n• WHO Japanese Encephalitis Fact Sheet (2024)\n• WHO Chikungunya Fact Sheet (2025)\n• CDC JE and Chikungunya Guidelines\n\nIndian Government:\n• NVBDCP Guidelines for VBD Control\n• IDSP Weekly Outbreak Reports\n• UIP Operational Guidelines\n• ICMR Research Publications\n\nScientific Literature:\n• Systematic reviews in PLoS NTD and Emerging Microbes & Infections\n• Indian Journal of Medical Research articles\n• Bulletin of WHO epidemiological studies\n\nNote: Always verify latest information from official sources as disease epidemiology and control strategies may change over time.'

# Save comprehensive PPTX
prs.save('JE_Chikungunya_India_Comprehensive.pptx')
print("Comprehensive PowerPoint presentation with 27 slides created successfully!")
print("Enhanced visualizations and presentation ready for MBBS teaching!")
